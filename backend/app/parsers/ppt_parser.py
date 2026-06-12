"""PowerPoint text extractor using python-pptx."""

from io import BytesIO

from pptx import Presentation

from app.parsers.base import BaseParser
from app.utils.logger import get_logger

logger = get_logger(__name__)


class PPTParser(BaseParser):
    """Extract text from .pptx files using python-pptx."""

    def parse(self, file_bytes: bytes, filename: str = "") -> str:
        """
        Iterate all slides and shapes, collecting text frame content.

        Args:
            file_bytes: Raw PPTX bytes.
            filename:   Original filename (informational).

        Returns:
            All slide text joined by newlines.
        """
        try:
            prs = Presentation(BytesIO(file_bytes))
            text_parts: list[str] = []
            for slide_num, slide in enumerate(prs.slides, start=1):
                slide_texts: list[str] = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            para_text = " ".join(
                                run.text for run in para.runs if run.text.strip()
                            )
                            if para_text.strip():
                                slide_texts.append(para_text)
                if slide_texts:
                    text_parts.append(f"[Slide {slide_num}]\n" + "\n".join(slide_texts))

            extracted = "\n\n".join(text_parts)
            logger.info(
                f"PPT parsed: {filename!r} — {len(prs.slides)} slides, "
                f"{len(extracted)} chars."
            )
            return extracted
        except Exception as exc:
            logger.error(f"PPTParser failed for {filename!r}: {exc}")
            raise
